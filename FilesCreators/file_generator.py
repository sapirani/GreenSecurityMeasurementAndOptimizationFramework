import numpy
from PIL import Image
from pathlib import Path
import random
import string
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor

DIR_NAME = "../Generated Files"
JPG_NAME = "jpg files"
POWERPOINT_NAME = "power point files"
TEXT_DIR_NAME = "text files"
DOC_DIR_NAME = "doc files"
NUMBER_OF_FILES = 10000

KB = 2**10
MB = 2**20


def generator(generate_file_func, dir_name, base_file_name, file_ending):
    """_summary_: genarate files in a directory

    Args:
        generate_file_func (func): function that generate the file
        dir_name (str): path to the directory
        base_file_name (str): base name of the file
        file_ending (str): file type
    """
    Path(f'{DIR_NAME}\\{dir_name}').mkdir(parents=True, exist_ok=True)
    for i in range(NUMBER_OF_FILES):
        generate_file_func(f'{DIR_NAME}\\{dir_name}\\{base_file_name}{i}.{file_ending}')

# generate all different types of files

def generate_text(full_dir):
    with open(f'{full_dir}', 'w') as f:
        f.write(rand_letters(random.randint(0.5 * MB, MB)))


def generate_jpg(full_dir):
    a = numpy.random.rand(random.randrange(750, 1500), random.randrange(750, 1500), 3) * 255
    im_out = Image.fromarray(a.astype('uint8')).convert('RGB')
    im_out.save(full_dir)


def generate_document(full_dir):
    document = Document()
    document.add_heading(rand_letters(20), 0)
    document.add_picture(get_random_jpg_file())
    document.add_paragraph(rand_letters(KB))
    document.save(full_dir)


def generate_powerpoint(full_dir):
    power_point = Presentation()

    # for i in range(20):
    slide = power_point.slides.add_slide(power_point.slide_layouts[0])
    slide.shapes.title.text = rand_letters(random.randint(15, 25))
    slide.placeholders[1].text = rand_letters(random.randint(100, 180))

    slide = power_point.slides.add_slide(power_point.slide_layouts[0])
    slide.shapes.add_picture(get_random_jpg_file(), 0, 0)

    power_point.save(full_dir)


def rand_letters(size):
    return ''.join([random.choice(string.ascii_letters) for i in range(size)])


def get_random_jpg_file():
    return f'{DIR_NAME}\\{JPG_NAME}\\pic{random.randint(0, NUMBER_OF_FILES - 1)}.jpg'


def main():
    generator(generate_text, TEXT_DIR_NAME, 'file', 'txt')
    #generator(generate_document, DOC_DIR_NAME, 'file', 'doc')
    #generator(generate_powerpoint, POWERPOINT_NAME, 'presentation', 'pptx')
    # generator(generate_jpg, JPG_NAME, 'pic', 'jpg')


if __name__ == '__main__':
    main()
